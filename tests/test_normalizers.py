"""Tests for brain.ingest.normalizers."""

from __future__ import annotations

from pathlib import Path
from unittest.mock import AsyncMock, patch

import pytest

from brain.ingest.normalizers import (
    get_normalizer,
    normalize_docx,
    normalize_pptx,
    normalize_text,
)
from brain.models import DocType


@pytest.fixture
def sample_txt(tmp_path):
    p = tmp_path / "sample.txt"
    p.write_text("Hello, this is a test document.\n\nWith two paragraphs.")
    return p


@pytest.fixture
def sample_docx(tmp_path):
    """Create a minimal .docx file using python-docx."""
    from docx import Document

    doc = Document()
    doc.add_paragraph("First paragraph of the docx.")
    doc.add_paragraph("Second paragraph of the docx.")
    p = tmp_path / "sample.docx"
    doc.save(str(p))
    return p


@pytest.fixture
def sample_pptx(tmp_path):
    """Create a minimal .pptx file using python-pptx."""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = "Slide 1 Title"
    slide.placeholders[1].text = "Slide 1 body text"
    p = tmp_path / "sample.pptx"
    prs.save(str(p))
    return p


def test_normalize_text(sample_txt):
    doc = normalize_text(sample_txt)
    assert doc.doc_type == DocType.TEXT
    assert "Hello" in doc.text
    assert "two paragraphs" in doc.text
    assert doc.content_hash != ""


def test_normalize_docx(sample_docx):
    doc = normalize_docx(sample_docx)
    assert doc.doc_type == DocType.DOCX
    assert "First paragraph" in doc.text
    assert "Second paragraph" in doc.text
    assert doc.content_hash != ""


def test_normalize_pptx(sample_pptx):
    doc = normalize_pptx(sample_pptx)
    assert doc.doc_type == DocType.PPTX
    assert "Slide 1" in doc.text
    assert "body text" in doc.text
    assert doc.metadata.get("slide_count") == 1


def test_get_normalizer_dispatch():
    assert get_normalizer(Path("test.docx")) is not None
    assert get_normalizer(Path("test.pptx")) is not None
    assert get_normalizer(Path("test.txt")) is not None
    assert get_normalizer(Path("test.md")) is not None
    assert get_normalizer(Path("test.png")) is not None
    assert get_normalizer(Path("test.unknown")) is None


@pytest.mark.asyncio
async def test_normalize_image_calls_llm(tmp_path):
    """Image normalizer should call transcribe_image from llm.py."""
    from brain.ingest.normalizers import normalize_image

    img = tmp_path / "test.png"
    img.write_bytes(b"\x89PNG\r\n\x1a\n" + b"\x00" * 100)

    with patch("brain.llm.transcribe_image", new_callable=AsyncMock) as mock:
        mock.return_value = "Transcribed content from image"
        doc = await normalize_image(img)

    assert doc.doc_type == DocType.IMAGE
    assert doc.text == "Transcribed content from image"
    mock.assert_called_once()
