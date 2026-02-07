"""Document normalizers — one function per file type.

Each normalizer reads a file and returns a Document with extracted text.
"""

from __future__ import annotations

import asyncio
from pathlib import Path

from brain.models import DocType, Document


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def normalize_docx(path: Path) -> Document:
    """Extract text from .docx preserving paragraph and table order."""
    from docx import Document as DocxDocument
    from docx.opc.constants import RELATIONSHIP_TYPE as RT
    from lxml import etree

    doc = DocxDocument(str(path))
    parts: list[str] = []

    # Walk the document body XML to interleave paragraphs and tables
    # in document order (python-docx iteration is paragraph-only or table-only).
    body = doc.element.body
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}

    for child in body:
        tag = etree.QName(child).localname
        if tag == "p":
            text = child.text or ""
            # Collect all runs
            runs = child.findall(".//w:t", ns)
            text = "".join(r.text or "" for r in runs)
            if text.strip():
                parts.append(text)
        elif tag == "tbl":
            rows = child.findall(".//w:tr", ns)
            for row in rows:
                cells = row.findall(".//w:tc", ns)
                cell_texts = []
                for cell in cells:
                    cell_runs = cell.findall(".//w:t", ns)
                    cell_texts.append("".join(r.text or "" for r in cell_runs))
                parts.append(" | ".join(cell_texts))

    title = doc.core_properties.title or path.stem
    full_text = "\n\n".join(parts)

    result = Document(
        source_path=str(path),
        doc_type=DocType.DOCX,
        title=title,
        text=full_text,
        metadata={"author": doc.core_properties.author or ""},
    )
    result.compute_hash()
    return result


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def normalize_pptx(path: Path) -> Document:
    """Extract text from .pptx — one section per slide with speaker notes."""
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_parts: list[str] = [f"--- Slide {slide_num} ---"]

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_parts.append(text)

            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cell_texts = [cell.text.strip() for cell in row.cells]
                    slide_parts.append(" | ".join(cell_texts))

        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_parts.append(f"[Speaker Notes] {notes}")

        parts.append("\n".join(slide_parts))

    title = path.stem
    if prs.core_properties and prs.core_properties.title:
        title = prs.core_properties.title

    full_text = "\n\n".join(parts)

    result = Document(
        source_path=str(path),
        doc_type=DocType.PPTX,
        title=title,
        text=full_text,
        metadata={"slide_count": len(prs.slides)},
    )
    result.compute_hash()
    return result


# ---------------------------------------------------------------------------
# Plain text (.txt, .md, .csv, .json, .yaml)
# ---------------------------------------------------------------------------

def normalize_text(path: Path) -> Document:
    """Read a plain text file."""
    text = path.read_text(encoding="utf-8", errors="replace")

    result = Document(
        source_path=str(path),
        doc_type=DocType.TEXT,
        title=path.stem,
        text=text,
    )
    result.compute_hash()
    return result


# ---------------------------------------------------------------------------
# Images (.png, .jpg, .jpeg, .webp, .heic)
# ---------------------------------------------------------------------------

async def normalize_image(path: Path) -> Document:
    """Transcribe image to text using LLM vision model."""
    from brain.llm import transcribe_image

    text = await transcribe_image(path)

    result = Document(
        source_path=str(path),
        doc_type=DocType.IMAGE,
        title=path.stem,
        text=text,
        metadata={"transcription_source": "llm_vision"},
    )
    result.compute_hash()
    return result


# ---------------------------------------------------------------------------
# Dispatch table
# ---------------------------------------------------------------------------

NORMALIZERS: dict[str, callable] = {
    ".docx": normalize_docx,
    ".pptx": normalize_pptx,
    ".txt": normalize_text,
    ".md": normalize_text,
    ".csv": normalize_text,
    ".json": normalize_text,
    ".yaml": normalize_text,
    ".yml": normalize_text,
    ".png": normalize_image,
    ".jpg": normalize_image,
    ".jpeg": normalize_image,
    ".webp": normalize_image,
    ".heic": normalize_image,
}


def get_normalizer(path: Path) -> callable | None:
    """Look up the normalizer for a file path, or None if unsupported."""
    return NORMALIZERS.get(path.suffix.lower())
